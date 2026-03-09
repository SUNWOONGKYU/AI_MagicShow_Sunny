# -*- coding: utf-8 -*-
"""슬라이드 5 오른쪽 컬럼 3분할
 - ① 특허출원 5건 (축소)
 - ② 풀스택 웹사이트 6개 제작 (신규 추가)
 - ③ 현재 3명 코칭 중 (사진 세로 확장)
"""
import sys, io, os, subprocess, glob, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

BASE   = "G:/내 드라이브/333_자료공유폴더/Sunny_Magic_Show"
PPTX   = BASE + "/MagicShow강의자료/AI_Magic_Show_v2_20260309.pptx"
LIBRE  = r"C:\Program Files\LibreOffice\program\soffice.exe"
TMPDIR = BASE + "/MagicShow강의자료/tmp_slide5"

def emu(inches): return int(inches * 914400)

prs = Presentation(PPTX)
slide5 = prs.slides[4]
shp = {s.shape_id: s for s in slide5.shapes}

# ── 레이아웃 계산 ──────────────────────────────────────────────────────────
# 오른쪽 컬럼 전체: t=0.950" ~ 5.425" (h=4.475")
# 3개 박스 + 2개 갭(0.100")
# Box1(특허)=1.050", Box2(풀스택)=1.050", Box3(코칭)=2.175"
# 갭=(4.475-1.050-1.050-2.175)/2 = 0.100"

COL_L = 4.700
COL_W = 4.571
INNER_L = 5.409   # 라벨/이미지 왼쪽
INNER_W = 3.141   # 라벨/이미지 폭

B1_TOP = 0.950;  B1_H = 1.050
B2_TOP = 2.100;  B2_H = 1.050
B3_TOP = 3.250;  B3_H = 2.175

LABEL_OFFSET = 0.075   # 박스 상단에서 라벨까지
LABEL_H      = 0.300
IMG_OFFSET   = 0.400   # 박스 상단에서 이미지까지
IMG_MARGIN   = 0.050   # 하단 여백

FILL_HEX = "003366"

# ── 1. Box 1 (특허출원 5건) 크기 조정 ─────────────────────────────────────
s11 = shp[11]
s11.top    = emu(B1_TOP)
s11.height = emu(B1_H)
print(f"  ✅ id=11 box1: t={B1_TOP}, h={B1_H}")

# 라벨 (id=12)
s12 = shp[12]
s12.top    = emu(B1_TOP + LABEL_OFFSET)
s12.height = emu(LABEL_H)
print(f"  ✅ id=12 label: t={B1_TOP + LABEL_OFFSET:.3f}")

# 특허 이미지 (id=13) — 크기 조정
s13 = shp[13]
s13.top    = emu(B1_TOP + IMG_OFFSET)
s13.height = emu(B1_H - IMG_OFFSET - IMG_MARGIN)   # = 0.600"
print(f"  ✅ id=13 patent img: t={B1_TOP + IMG_OFFSET:.3f}, h={B1_H - IMG_OFFSET - IMG_MARGIN:.3f}")

# ── 2. Box 2 (풀스택 웹사이트 6개 제작) 신규 추가 ─────────────────────────
# 배경 박스
bg2 = slide5.shapes.add_shape(1, emu(COL_L), emu(B2_TOP), emu(COL_W), emu(B2_H))
bg2.fill.solid()
bg2.fill.fore_color.rgb = RGBColor.from_string(FILL_HEX)
bg2.line.fill.background()
print(f"  ✅ Box2 배경 추가")

# 라벨
lbl2 = slide5.shapes.add_textbox(emu(INNER_L), emu(B2_TOP + LABEL_OFFSET),
                                   emu(INNER_W + 0.500), emu(LABEL_H))
tf = lbl2.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
run = p.add_run()
run.text = "🌐 풀스택 웹사이트 6개 제작"
run.font.size = Pt(14)
run.font.bold = True
run.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00)
print(f"  ✅ Box2 라벨 추가")

# 웹사이트 목록 텍스트
SITES = [
    "SSAL Works  /  Politician Finder  /  ValueLink",
    "MyChatbot World  /  AX-On  /  FX Pooling",
]
txt2 = slide5.shapes.add_textbox(emu(INNER_L - 0.200), emu(B2_TOP + 0.400),
                                   emu(INNER_W + 0.700), emu(0.580))
tf2 = txt2.text_frame
tf2.word_wrap = True
for i, line in enumerate(SITES):
    p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    r2 = p2.add_run()
    r2.text = line
    r2.font.size = Pt(10)
    r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
print(f"  ✅ Box2 목록 텍스트 추가")

# ── 3. Box 3 (현재 3명 코칭 중) 이동 + 사진 확장 ──────────────────────────
s20 = shp[20]
s20.top    = emu(B3_TOP)
s20.height = emu(B3_H)
print(f"  ✅ id=20 box3: t={B3_TOP}, h={B3_H}")

# 라벨 (id=21)
s21 = shp[21]
s21.top    = emu(B3_TOP + LABEL_OFFSET)
s21.height = emu(LABEL_H)
print(f"  ✅ id=21 label: t={B3_TOP + LABEL_OFFSET:.3f}")

# 코칭 사진 (id=22) — 세로 확장 (찌그러짐 수정: 원래 1.572 → 1.725")
s22 = shp[22]
s22.top    = emu(B3_TOP + IMG_OFFSET)
new_photo_h = B3_H - IMG_OFFSET - IMG_MARGIN   # = 1.725"
s22.height = emu(new_photo_h)
print(f"  ✅ id=22 coaching photo: t={B3_TOP + IMG_OFFSET:.3f}, h={new_photo_h:.3f}")

# ── 저장 ─────────────────────────────────────────────────────────────────
prs.save(PPTX)
print(f"\n✅ PPTX 저장 완료")

# ── 슬라이드 5 PNG 재생성 ─────────────────────────────────────────────────
os.makedirs(TMPDIR, exist_ok=True)
result = subprocess.run([
    LIBRE, "--headless", "--convert-to", "pdf",
    "--outdir", TMPDIR, PPTX
], capture_output=True, text=True, timeout=120)

pdf_files = glob.glob(TMPDIR + "/*.pdf")
if not pdf_files:
    print("ERR: PDF 생성 실패", result.stderr[:200])
else:
    import fitz
    doc = fitz.open(pdf_files[0])
    page = doc[4]   # 0-indexed, slide 5
    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
    out_png = BASE + "/MagicShow강의자료/slides/slide5.png"
    pix.save(out_png)
    doc.close()
    print(f"✅ slide5.png 교체 완료")

shutil.rmtree(TMPDIR)
print("✅ 완료")
