# -*- coding: utf-8 -*-
"""슬라이드 7 카드 간격/폰트 수정 v2
 - 모든 이름 bold=True 통일 (SSAL, Politician, MyChatbot 비bold 수정)
 - 이름 top offset: 0.100" → 0.150" (컬러바와 간격 0.020"→0.070")
 - desc top 유지 → 이름-desc 간격: 0.080" → 0.030" (좁힘)
 - desc 폰트 11pt 통일 (AX-On/FX Pooling 이미 11pt, 나머지 12pt→11pt)
"""
import sys, io, os, subprocess, glob, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree
from pptx.oxml.ns import qn

BASE   = "G:/내 드라이브/333_자료공유폴더/Sunny_Magic_Show"
PPTX   = BASE + "/MagicShow강의자료/AI_Magic_Show_v2_20260309.pptx"
LIBRE  = r"C:\Program Files\LibreOffice\program\soffice.exe"
TMPDIR = BASE + "/MagicShow강의자료/tmp_page7v2"

def emu(inches): return int(inches * 914400)

prs = Presentation(PPTX)
slide7 = prs.slides[6]
shp = {s.shape_id: s for s in slide7.shapes}

# ── 1. 이름 shape: bold=True 통일 + top 이동 ──────────────────────────────
#   Card tops: row1=0.950", row2=3.300"
#   Color bar height: 0.080"
#   New name offset: 0.150" from card top (was 0.100")
NAME_SHAPE_INFO = {
    # shape_id: (card_top_inch,)
    7:  0.950,  # SSAL
    13: 0.950,  # Politician
    19: 0.950,  # ValueLink
    25: 3.300,  # MyChatbot
    31: 3.300,  # AX-On
    37: 3.300,  # FX Pooling
}
NEW_NAME_OFFSET = 0.150   # from card_top

for sid, card_top in NAME_SHAPE_INFO.items():
    s = shp.get(sid)
    if not s:
        print(f"  ⚠ id={sid} 없음")
        continue
    # 위치 이동
    new_top = emu(card_top + NEW_NAME_OFFSET)
    s.top = new_top
    # 모든 run → bold=True, 16pt
    for para in s.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                rPr = etree.SubElement(run._r, qn('a:rPr'))
            rPr.set('b', '1')
            rPr.set('sz', '1600')
    print(f"  ✅ name id={sid}: top={card_top + NEW_NAME_OFFSET:.3f}\" bold=True")

# ── 2. desc shape: 폰트 11pt 통일 ─────────────────────────────────────────
DESC_IDS = [8, 14, 20, 26, 32, 38]

for sid in DESC_IDS:
    s = shp.get(sid)
    if not s:
        continue
    changed = 0
    for para in s.text_frame.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                rPr = etree.SubElement(run._r, qn('a:rPr'))
            cur_sz = rPr.get('sz')
            if cur_sz != '1100':
                rPr.set('sz', '1100')
                changed += 1
    print(f"  ✅ desc id={sid}: {changed}개 run → 11pt")

# ── 저장 ─────────────────────────────────────────────────────────────────
prs.save(PPTX)
print(f"\n✅ PPTX 저장 완료")

# ── 슬라이드 7 PNG 재생성 ──────────────────────────────────────────────────
os.makedirs(TMPDIR, exist_ok=True)
result = subprocess.run([
    LIBRE, "--headless", "--convert-to", "pdf",
    "--outdir", TMPDIR, PPTX
], capture_output=True, text=True, timeout=120)

pdf_files = glob.glob(TMPDIR + "/*.pdf")
if not pdf_files:
    print("ERR: PDF 생성 실패", result.stderr[:200])
else:
    import fitz
    doc = fitz.open(pdf_files[0])
    page = doc[6]
    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
    out_png = BASE + "/MagicShow강의자료/slides/slide7.png"
    pix.save(out_png)
    doc.close()
    print(f"✅ slide7.png 교체 완료")

shutil.rmtree(TMPDIR)
print("✅ 임시 파일 정리 완료")
