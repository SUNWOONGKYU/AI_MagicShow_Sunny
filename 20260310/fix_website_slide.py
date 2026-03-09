# -*- coding: utf-8 -*-
"""슬라이드 7 (웹사이트 6개) 수정"""
import sys, io, os, subprocess, glob, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

BASE    = "G:/내 드라이브/333_자료공유폴더/Sunny_Magic_Show"
PPTX    = BASE + "/MagicShow강의자료/AI_Magic_Show_v2_20260309.pptx"
LIBRE   = r"C:\Program Files\LibreOffice\program\soffice.exe"
TMPDIR  = BASE + "/MagicShow강의자료/tmp_page7"

def emu(inches): return int(inches * 914400)

prs = Presentation(PPTX)
slide7 = prs.slides[6]  # 0-indexed → 슬라이드 7

shape_by_id = {s.shape_id: s for s in slide7.shapes}

# ── 1. 링크 관련 shape 전부 삭제 ─────────────────────────────────────
# 링크 배경(bg): 9,15,21,27,33,39  /  링크 텍스트: 10,16,22,28,34,40
LINK_SHAPE_IDS = [9, 10, 15, 16, 21, 22, 27, 28, 33, 34, 39, 40]
for sid in LINK_SHAPE_IDS:
    s = shape_by_id.get(sid)
    if s:
        s._element.getparent().remove(s._element)
        print(f"  ✅ shape_id={sid} 삭제")

# ── 2. 각 카드 desc 높이 확장 (링크 빠진 공간 채우기) ────────────────
# DESC 시작 = card_top + 0.680"  →  새 높이 = 1.370" (card bottom - 0.100")
# card tops: row1=0.950", row2=3.300"
CARD_DESC_IDS = [8, 14, 20, 26, 32, 38]  # desc shape ids
NEW_DESC_H = emu(1.370)
for sid in CARD_DESC_IDS:
    s = shape_by_id.get(sid)
    if s:
        s.height = NEW_DESC_H
        print(f"  ✅ shape_id={sid} desc 높이 확장")

# ── 3. ValueLine → ValueLink (두 줄) ─────────────────────────────────
name_shape = shape_by_id.get(19)  # Card 3 (ValueLine) 이름 shape
if name_shape:
    tf = name_shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if 'ValueLine' in run.text or 'Value' in run.text:
                run.text = 'Value'
                # 두 번째 줄 추가
                from pptx.oxml.ns import qn
                from lxml import etree
                # 새 paragraph 추가
                p2 = tf.add_paragraph()
                p2.alignment = para.alignment
                r2 = p2.add_run()
                r2.text = 'Link'
                r2.font.size = Pt(16)
                r2.font.bold = True
                r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                print(f"  ✅ ValueLine → ValueLink (두 줄)")
                break

# ── 4. AX-On/FX Pooling: bold 제거, 폰트 크기 맞춤 (16pt = 203200) ──
# 이름 shape: 31(AX-On), 37(FX Pooling)  / desc shape: 32, 38
for sid in [31, 37]:
    s = shape_by_id.get(sid)
    if s:
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True   # 다른 카드와 동일하게 bold
                run.font.size = Pt(16)
        print(f"  ✅ shape_id={sid} 이름 폰트 정렬 (16pt, bold=True)")

# desc 교체 + 폰트 크기 맞춤
DESC_UPDATES = {
    32: "분야별 AI 전문가들이 모여서\n시니어 AI 창업 교육과\n기업 AX 프로젝트를 추진하는 플랫폼",   # AX-On
    38: "외환풀링 및\n스테이블코인을 활용한\n실시간 환위험 관리 시스템",               # FX Pooling
}
for sid, new_text in DESC_UPDATES.items():
    s = shape_by_id.get(sid)
    if s:
        tf = s.text_frame
        # 기존 단락 모두 제거 후 재구성
        from pptx.oxml.ns import qn
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
        # 첫 단락 재사용
        p0 = tf.paragraphs[0]
        for run in list(p0.runs):
            run._r.getparent().remove(run._r)
        lines = new_text.split('\n')
        r0 = p0.add_run()
        r0.text = lines[0]
        r0.font.size = Pt(11)
        r0.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        for line in lines[1:]:
            pn = tf.add_paragraph()
            pn.alignment = p0.alignment
            rn = pn.add_run()
            rn.text = line
            rn.font.size = Pt(11)
            rn.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        print(f"  ✅ shape_id={sid} desc 교체 완료")

# ── 저장 ─────────────────────────────────────────────────────────────
prs.save(PPTX)
print(f"\n✅ PPTX 저장 완료")

# ── 슬라이드 7만 PNG 재생성 ───────────────────────────────────────────
os.makedirs(TMPDIR, exist_ok=True)
result = subprocess.run([
    LIBRE, "--headless", "--convert-to", "pdf",
    "--outdir", TMPDIR, PPTX
], capture_output=True, text=True, timeout=120)

pdf_files = glob.glob(TMPDIR + "/*.pdf")
if not pdf_files:
    print("ERR: PDF 생성 실패", result.stderr[:200])
else:
    import fitz
    doc = fitz.open(pdf_files[0])
    page = doc[6]   # 0-indexed → slide 7
    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
    out_png = BASE + "/MagicShow강의자료/slides/slide7.png"
    pix.save(out_png)
    doc.close()
    print(f"✅ slide7.png 교체 완료")

shutil.rmtree(TMPDIR)
print("✅ 임시 파일 정리 완료")
