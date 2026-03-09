# -*- coding: utf-8 -*-
"""AI Magic Show 강의자료 PPTX 수정 스크립트"""
import sys, io, copy, shutil, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── 경로 설정 ──────────────────────────────────────────────────────────
BASE = "G:/내 드라이브/333_자료공유폴더/Sunny_Magic_Show"
PPTX_IN  = BASE + "/MagicShow강의자료/AI_Magic_Show_53.pptx"
PPTX_OUT = BASE + "/MagicShow강의자료/AI_Magic_Show_수정본.pptx"
PATENT_DIR = BASE + "/patent"

PATENT_IMAGES = [
    (BASE + "/patent/스테이블코인.png",            "10-2025-0090930\n원화 스테이블코인 기반\n온라인 커뮤니티 경제\n생태계 플랫폼"),
    (BASE + "/patent/AI아바타.png",               "10-2025-0095665\n대규모 객체형 AI 아바타\n멀티 플랫폼 통합\n서비스"),
    (BASE + "/patent/Patent_SALGrid/특허출원번호통지서.jpg",     "10-2026-0009425\nSAL Grid\n3차원 좌표 기반\n멀티태스크 오케스트레이션"),
    (BASE + "/patent/Patent_MyChatbot/챗봇특허출원번호통지서.png","10-2026-0038658\nMyChatbot\n멀티 페르소나\nAI 챗봇 라이프사이클"),
    (BASE + "/patent/Platoon_Formation_Patent/특허출원번호통지서.png","10-2026-0041235\n소대편제\n계층적 다중 AI\n에이전트 작업팀"),
]

# Emu 변환 헬퍼
def emu(inches): return int(inches * 914400)

# ── EMU 상수 ────────────────────────────────────────────────────────────
SLIDE_W = emu(10.0)
SLIDE_H = emu(5.625)

# 새 카드 레이아웃 (3열 × 2행)
CARD_W = emu(3.100)
CARD_H = emu(2.150)
COL_GAP = emu(0.150)
ROW_GAP = emu(0.200)
LEFT_MARGIN = emu(0.200)
ROW1_TOP = emu(0.950)
ROW2_TOP = emu(0.950 + 2.150 + 0.200)  # = 3.300"

COL_LEFTS = [
    LEFT_MARGIN,
    LEFT_MARGIN + CARD_W + COL_GAP,
    LEFT_MARGIN + CARD_W + COL_GAP + CARD_W + COL_GAP,
]

# 카드 내부 요소 (card 기준 오프셋)
INSET_X = emu(0.100)
BAR_H   = emu(0.080)
NAME_Y  = emu(0.100);  NAME_H  = emu(0.500)
DESC_Y  = emu(0.680);  DESC_H  = emu(0.680)
LINK_Y  = emu(1.600);  LINK_H  = emu(0.400)
LINK_W_EXTRA = emu(0.050)  # link text slightly wider

INNER_W = CARD_W - 2 * INSET_X

# 카드 데이터 (Row1: 0,1,2 / Row2: 3,4,5)
CARDS = [
    # Row 1
    {"name": "SSAL\nWorks",      "desc": "AI를 활용한/\nSAL Grid 기반\n웹사이트 개발 플랫폼", "url": "🔗 ssalworks.ai.kr",            "bar": "2ECC71"},
    {"name": "Politician\nFinder","desc": "AI 기반의 정치인\n평가/커뮤니티\n플랫폼",           "url": "🔗 politicianfinder.ai.kr",       "bar": "3498DB"},
    {"name": "ValueLine",         "desc": "AI를 활용한\n기업가치\n평가/연결 플랫폼",            "url": "🔗 sunwoongkyu.github.io/ValueLink","bar": "F39C12"},
    # Row 2
    {"name": "MyChatbot\nWorld",  "desc": "AI 챗봇\n라이프사이클\n플랫폼",                     "url": "🔗 mychatbot.world",              "bar": "9B59B6"},
    {"name": "AX-On\nPlatform",   "desc": "AI 전환(AX)\n전문가·창업·프로젝트\n연결 플랫폼",  "url": "🔗 ax-on-platform.vercel.app",    "bar": "E74C3C"},
    {"name": "FX Pooling\nPlatform","desc":"외환 풀링\n전문 플랫폼",                          "url": "🔗 prototype-six-coral.vercel.app","bar": "1ABC9C"},
]

def card_pos(card_idx):
    """카드 인덱스(0-5) → (left, top) emu"""
    row = card_idx // 3
    col = card_idx % 3
    left = COL_LEFTS[col]
    top  = ROW1_TOP if row == 0 else ROW2_TOP
    return left, top


# ── 헬퍼 함수들 ──────────────────────────────────────────────────────────

def set_shape_pos(shape, left, top, width, height):
    shape.left   = left
    shape.top    = top
    shape.width  = width
    shape.height = height

def set_fill_solid(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)

def set_text_color(shape, hex_color):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor.from_string(hex_color)

def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color="FFFFFF", align=PP_ALIGN.LEFT, wrap=True):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    # 줄바꿈 포함 처리
    if '\n' in text:
        lines = text.split('\n')
        run.text = lines[0]
        for line in lines[1:]:
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            from lxml import etree
            from pptx.oxml.ns import qn
            br = etree.SubElement(p._p, qn('a:br'))
            p2 = tf.add_paragraph()
            p2.alignment = align
            r2 = p2.add_run()
            r2.text = line
            r2.font.size = Pt(font_size)
            r2.font.bold = bold
            r2.font.color.rgb = RGBColor.from_string(color)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    set_fill_solid(shape, fill_color)
    if line_color:
        shape.line.color.rgb = RGBColor.from_string(line_color)
    else:
        shape.line.fill.background()  # no line
    return shape


def delete_slide(prs, slide_idx):
    """슬라이드 삭제"""
    from pptx.oxml.ns import nsmap
    xml_slides = prs.slides._sldIdLst
    # r:id 속성 가져오기 (namespace 포함)
    sldId_elem = xml_slides[slide_idx]
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    rId = sldId_elem.get(f'{{{r_ns}}}id')
    # 관계 제거
    prs.part.drop_rel(rId)
    # sldId 요소 제거
    xml_slides.remove(sldId_elem)


def duplicate_slide(prs, slide_idx):
    """슬라이드 복제 후 마지막에 추가 (나중에 이동)"""
    template = prs.slides[slide_idx]
    blank_layout = prs.slide_layouts[6]  # blank

    # slide XML 복사
    from pptx.opc.packuri import PackURI
    slide_part = template.part
    # 새 슬라이드 생성
    new_slide = prs.slides.add_slide(blank_layout)
    # 복사한 XML로 교체
    new_slide.shapes._spTree.clear()
    for elem in template.shapes._spTree:
        new_slide.shapes._spTree.append(copy.deepcopy(elem))
    return new_slide


def move_slide_to(prs, from_idx, to_idx):
    """슬라이드를 from_idx → to_idx 위치로 이동"""
    xml_slides = prs.slides._sldIdLst
    slide_elem = xml_slides[from_idx]
    xml_slides.remove(slide_elem)
    xml_slides.insert(to_idx, slide_elem)


# ── Main ─────────────────────────────────────────────────────────────────
print("▶ PPTX 로드 중...")
prs = Presentation(PPTX_IN)
print(f"  총 슬라이드 수: {len(prs.slides)}")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# [1] Slide 5: "특허출원 3건" → "특허출원 5건"
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("\n[1] Slide 5: 특허출원 3건 → 5건")
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if hasattr(shape, "text") and "특허출원 3건" in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "3건" in run.text:
                    run.text = run.text.replace("3건", "5건")
                    print(f"  ✅ shape_id={shape.shape_id}: '3건' → '5건'")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# [2] Slide 6: 4개 → 6개 레이아웃 재구성
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("\n[2] Slide 6: 웹사이트 4개 → 6개 (3×2 레이아웃)")
slide6 = prs.slides[5]

# 제목 변경
for shape in slide6.shapes:
    if hasattr(shape, "text") and "4개 제작" in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "4개" in run.text:
                    run.text = run.text.replace("4개", "6개")
                    print(f"  ✅ 제목: '4개' → '6개'")

# 기존 카드 shape 매핑 (shape_id 기준)
# Card 1: container=5, topbar=6, name=7, desc=8, linkbg=9, linktext=10
# Card 2: container=11, topbar=12, name=13, desc=14, linkbg=15, linktext=16
# Card 3: container=17, topbar=18, name=19, desc=20, linkbg=21, linktext=22
# Card 4: container=23, topbar=24, name=25, desc=26, linkbg=27, linktext=28
CARD_SHAPE_IDS = [
    [5, 6, 7, 8, 9, 10],   # Card 1 SSAL
    [11,12,13,14,15,16],   # Card 2 Politician
    [17,18,19,20,21,22],   # Card 3 ValueLine
    [23,24,25,26,27,28],   # Card 4 MyChatbot
]
shape_by_id = {s.shape_id: s for s in slide6.shapes}

for card_idx, ids in enumerate(CARD_SHAPE_IDS):
    left, top = card_pos(card_idx)
    cdata = CARDS[card_idx]

    sid_container, sid_bar, sid_name, sid_desc, sid_linkbg, sid_linktxt = ids

    # Container
    s = shape_by_id[sid_container]
    set_shape_pos(s, left, top, CARD_W, CARD_H)

    # Top bar
    s = shape_by_id[sid_bar]
    set_shape_pos(s, left, top, CARD_W, BAR_H)
    set_fill_solid(s, cdata["bar"])

    # Name
    s = shape_by_id[sid_name]
    set_shape_pos(s, left + INSET_X, top + NAME_Y, INNER_W, NAME_H)

    # Desc
    s = shape_by_id[sid_desc]
    set_shape_pos(s, left + INSET_X, top + DESC_Y, INNER_W, DESC_H)

    # Link bg
    s = shape_by_id[sid_linkbg]
    set_shape_pos(s, left + INSET_X, top + LINK_Y, INNER_W, LINK_H)

    # Link text
    s = shape_by_id[sid_linktxt]
    set_shape_pos(s, left + INSET_X - LINK_W_EXTRA, top + LINK_Y + emu(0.010),
                  INNER_W + LINK_W_EXTRA * 2, LINK_H - emu(0.020))
    # URL 텍스트 업데이트
    for para in s.text_frame.paragraphs:
        for run in para.runs:
            run.text = cdata["url"]

    print(f"  ✅ Card {card_idx+1} ({cdata['name'].split(chr(10))[0]}) 재배치 완료")

# 새 카드 2개 추가 (Card 5, 6) — Card 4의 XML을 복사해서 수정
print("  ➕ 새 카드 5, 6 추가 중...")

from pptx.util import Pt
from pptx.oxml.ns import qn as nsqn

def add_card_to_slide(slide, card_idx, cdata):
    """새 카드를 슬라이드에 추가"""
    left, top = card_pos(card_idx)

    # 1. Container background
    c = add_rect(slide, left, top, CARD_W, CARD_H, "003366")

    # 2. Top color bar
    b = add_rect(slide, left, top, CARD_W, BAR_H, cdata["bar"])

    # 3. Name text
    n = slide.shapes.add_textbox(left + INSET_X, top + NAME_Y, INNER_W, NAME_H)
    n.text_frame.word_wrap = True
    p = n.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cdata["name"]
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 4. Description text
    d = slide.shapes.add_textbox(left + INSET_X, top + DESC_Y, INNER_W, DESC_H)
    d.text_frame.word_wrap = True
    p = d.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cdata["desc"]
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 5. Link button background
    lb = add_rect(slide, left + INSET_X, top + LINK_Y, INNER_W, LINK_H, "1A3A6A")

    # 6. Link text
    lt = slide.shapes.add_textbox(left + INSET_X - LINK_W_EXTRA,
                                   top + LINK_Y + emu(0.010),
                                   INNER_W + LINK_W_EXTRA * 2,
                                   LINK_H - emu(0.020))
    lt.text_frame.word_wrap = False
    p = lt.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cdata["url"]
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x85, 0xC1, 0xE9)

for ci in [4, 5]:
    add_card_to_slide(slide6, ci, CARDS[ci])
    print(f"  ✅ Card {ci+1} ({CARDS[ci]['name'].split(chr(10))[0]}) 추가 완료")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# [3] 특허출원번호통지서 슬라이드 신규 추가 (Slide 5 뒤에)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("\n[3] 특허출원번호통지서 5건 슬라이드 신규 추가")

# blank 레이아웃으로 새 슬라이드 만들기 (마지막에 추가됨)
blank_layout = prs.slide_layouts[0]  # DEFAULT (유일한 레이아웃)

new_patent_slide = prs.slides.add_slide(blank_layout)

# 어두운 배경 추가 (전체 슬라이드 덮는 사각형)
bg = new_patent_slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x00, 0x10, 0x20)
bg.line.fill.background()

# 주황색 구분선 (기존 슬라이드 스타일)
line = new_patent_slide.shapes.add_shape(1, emu(0.4), emu(0.72), emu(9.2), emu(0.02))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0xFF, 0x8C, 0x00)
line.line.fill.background()

# 제목
title_box = new_patent_slide.shapes.add_textbox(emu(0.4), emu(0.20), emu(9.2), emu(0.45))
title_tf = title_box.text_frame
p = title_tf.paragraphs[0]
run = p.add_run()
run.text = "📋  특허출원번호통지서  5건"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00)

# 5개 이미지 레이아웃 (3 + 2)
# 사용 가능 영역: y=0.85" ~ y=5.40" (h=4.55"), x=0.20" ~ x=9.80" (w=9.60")
IMG_AREA_TOP = emu(0.85)
IMG_AREA_H   = emu(4.55)
IMG_AREA_LEFT = emu(0.20)
IMG_AREA_W    = emu(9.60)
IMG_GAP_X = emu(0.20)
IMG_GAP_Y = emu(0.15)

# Row 1: 3개 (idx 0-2)
row1_h = (IMG_AREA_H - IMG_GAP_Y) * 0.50  # 두 행 동일 높이
row1_w = (IMG_AREA_W - 2 * IMG_GAP_X) / 3
row1_top = IMG_AREA_TOP

row2_h = IMG_AREA_H - row1_h - IMG_GAP_Y
row2_w = (IMG_AREA_W - IMG_GAP_X) / 2
row2_top = IMG_AREA_TOP + row1_h + IMG_GAP_Y

# row2 2개를 가운데 정렬 (전체 너비에서)
row2_total_w = row2_w * 2 + IMG_GAP_X
row2_left_start = IMG_AREA_LEFT + (IMG_AREA_W - row2_total_w) / 2

positions = [
    # Row 1 (3개)
    (IMG_AREA_LEFT,                          row1_top, row1_w, row1_h),
    (IMG_AREA_LEFT + row1_w + IMG_GAP_X,     row1_top, row1_w, row1_h),
    (IMG_AREA_LEFT + 2*(row1_w + IMG_GAP_X), row1_top, row1_w, row1_h),
    # Row 2 (2개, 가운데 정렬)
    (row2_left_start,              row2_top, row2_w, row2_h),
    (row2_left_start + row2_w + IMG_GAP_X, row2_top, row2_w, row2_h),
]

for i, (img_path, label) in enumerate(PATENT_IMAGES):
    l, t, w, h = [int(x) for x in positions[i]]

    # 이미지 배경 박스 (흰 테두리)
    frame = new_patent_slide.shapes.add_shape(1, l, t, w, h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    frame.line.color.rgb = RGBColor(0xFF, 0x8C, 0x00)
    frame.line.width = Pt(2)

    # 이미지 삽입 (5% 안쪽 여백)
    pad = int(w * 0.02)
    if os.path.exists(img_path):
        new_patent_slide.shapes.add_picture(
            img_path, l + pad, t + pad, w - 2*pad, h - 2*pad
        )
        print(f"  ✅ 이미지 {i+1}: {os.path.basename(img_path)}")
    else:
        # 이미지 없을 때 텍스트로 대체
        tb = new_patent_slide.shapes.add_textbox(l+pad, t+pad, w-2*pad, h-2*pad)
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"[이미지 없음]\n{label}"
        run.font.color.rgb = RGBColor(0xFF, 0, 0)
        print(f"  ⚠️ 이미지 없음: {img_path}")

# 새 슬라이드를 Slide 5 (index 4) 다음으로 이동
# 현재 마지막 위치 → index 4+1=5로 이동
last_idx = len(prs.slides) - 1
move_slide_to(prs, last_idx, 5)
print(f"  ✅ 특허 통지서 슬라이드 → Slide 6 위치로 이동")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# [4] Vanilla vs React 슬라이드 삭제
# 새 슬라이드 삽입 후: 기존 Slide 43 → Slide 44 (index 43)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("\n[4] Vanilla vs React 슬라이드 삭제")
# 삭제 전 확인
vanilla_idx = None
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text") and "Vanilla vs React" in shape.text:
            vanilla_idx = idx
            break
    if vanilla_idx is not None:
        break

if vanilla_idx is not None:
    print(f"  발견: Slide {vanilla_idx+1} (index {vanilla_idx})")
    delete_slide(prs, vanilla_idx)
    print(f"  ✅ 삭제 완료")
else:
    print("  ⚠️ 슬라이드를 찾지 못함")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# [5] Tips 번호 재조정: Tip 4→3, Tip 5→4, ..., Tip 10→9
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("\n[5] Tips 번호 재조정 (Tip 4→3, 5→4, ..., 10→9)")
# 연쇄 치환 방지: run 단위로 한 번만 적용 (높은 번호부터)
for slide in prs.slides:
    slide_idx = list(prs.slides).index(slide)
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for num in range(10, 3, -1):  # 10 → 4 순서로 확인
                    old = f"Tip {num}:"
                    new = f"Tip {num-1}:"
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        print(f"  ✅ '{old}' → '{new}' (Slide {slide_idx+1})")
                        break  # 이 run에서 한 번만 적용


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 저장
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print(f"\n▶ 저장 중: {PPTX_OUT}")
prs.save(PPTX_OUT)
print(f"✅ 완료! 총 슬라이드: {len(prs.slides)}")
