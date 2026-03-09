# -*- coding: utf-8 -*-
"""특허출원번호통지서 슬라이드 수정 — 이미지 비율 유지 + 제목 변경"""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

BASE = "G:/내 드라이브/333_자료공유폴더/Sunny_Magic_Show"
PPTX_IN  = BASE + "/MagicShow강의자료/AI_Magic_Show_v2_20260309.pptx"
PPTX_OUT = BASE + "/MagicShow강의자료/AI_Magic_Show_v2_20260309.pptx"  # 덮어씀

PATENT_IMAGES = [
    # Row 1 (상단 3개 — 가로 좁은 것)
    BASE + "/patent/Patent_SALGrid/특허출원번호통지서.jpg",
    BASE + "/patent/Patent_MyChatbot/챗봇특허출원번호통지서.png",
    BASE + "/patent/Platoon_Formation_Patent/특허출원번호통지서.png",
    # Row 2 (하단 2개 — 가로 넓은 것)
    BASE + "/patent/스테이블코인.png",
    BASE + "/patent/AI아바타.png",
]

def emu(inches): return int(inches * 914400)

SLIDE_W = emu(10.0)
SLIDE_H = emu(5.625)

# 각 이미지의 실제 가로/세로 비율 계산
ratios = []
for p in PATENT_IMAGES:
    img = Image.open(p)
    w, h = img.size
    ratios.append(w / h)
    print(f"  {os.path.basename(p)}: {w}x{h} (ratio={w/h:.3f})")

# ── 레이아웃 계산 (비율 유지, 동일 높이) ──────────────────────────────
# 사용 영역: y=0.82" ~ y=5.55" (h=4.73")
# 행 간격: 0.15", 좌우 마진: 0.20"
ROW_H_IN = 2.28   # 행 높이 (인치) — 두 행 동일
ROW1_TOP  = 0.82
ROW2_TOP  = ROW1_TOP + ROW_H_IN + 0.17  # = 3.27"
H_GAP = 0.18      # 이미지 사이 수평 간격

# Row 1: 이미지 0, 1, 2 → 각 width = ratio * ROW_H_IN
row1_widths = [ratios[i] * ROW_H_IN for i in range(3)]
row1_total_w = sum(row1_widths) + H_GAP * 2
row1_left_start = (10.0 - row1_total_w) / 2  # 가운데 정렬

# Row 2: 이미지 3, 4
row2_widths = [ratios[i] * ROW_H_IN for i in range(3, 5)]
row2_total_w = sum(row2_widths) + H_GAP
row2_left_start = (10.0 - row2_total_w) / 2  # 가운데 정렬

print(f"\nRow 1 widths: {[f'{w:.2f}' for w in row1_widths]}, total={row1_total_w:.2f}\"")
print(f"Row 2 widths: {[f'{w:.2f}' for w in row2_widths]}, total={row2_total_w:.2f}\"")
print(f"Row 1 starts x={row1_left_start:.2f}\"")
print(f"Row 2 starts x={row2_left_start:.2f}\"")

# 이미지별 (left, top, width, height) 인치 단위
positions_in = []
# Row 1
x = row1_left_start
for i in range(3):
    positions_in.append((x, ROW1_TOP, row1_widths[i], ROW_H_IN))
    x += row1_widths[i] + H_GAP
# Row 2
x = row2_left_start
for i in range(2):
    positions_in.append((x, ROW2_TOP, row2_widths[i], ROW_H_IN))
    x += row2_widths[i] + H_GAP

# ── PPTX 수정 ──────────────────────────────────────────────────────────
print("\n▶ PPTX 로드...")
prs = Presentation(PPTX_IN)

# Slide index 5 = 새로 만든 특허통지서 슬라이드
patent_slide = prs.slides[5]

# 기존 shapes 전부 삭제 (배경 사각형 제외하고 싶지만 전체 재구성이 더 깔끔)
sp_tree = patent_slide.shapes._spTree
# 배경·구분선·제목 포함 모든 shape 제거 후 재구성
for shape in list(patent_slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

print("  기존 shapes 삭제 완료")

# ── 슬라이드 배경색 설정 (PPTX 자체 배경 — PowerPoint에서도 보임) ────
bg = patent_slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x00, 0x10, 0x20)
print("  슬라이드 배경색 설정 완료 (#001020)")

# ── 배경 재구성 ──────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_hex, line=False):
    s = slide.shapes.add_shape(1, emu(l), emu(t), emu(w), emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if not line:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGBColor.from_string("FF8C00")
        s.line.width = Pt(1.5)
    return s

# 배경
add_rect(patent_slide, 0, 0, 10.0, 5.625, "001020")
# 주황 구분선 (슬라이드 7 기준: top=0.72")
add_rect(patent_slide, 0.40, 0.72, 9.20, 0.02, "FF8C00")

# 제목 — 슬라이드 7과 동일 형식 (left=0.95", top=0.20", 24pt, bold, 주황)
tb = patent_slide.shapes.add_textbox(emu(0.95), emu(0.20), emu(8.50), emu(0.45))
tf = tb.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
run = p.add_run()
run.text = "📋  특허출원  5건"
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = RGBColor(0xFF, 0x8C, 0x00)

print("  배경·제목 추가 완료")

# ── 이미지 삽입 (비율 유지) ──────────────────────────────────────────
PAD = 0.03  # 프레임 안쪽 여백 (인치)
for i, (img_path, (l, t, w, h)) in enumerate(zip(PATENT_IMAGES, positions_in)):
    # 프레임 (흰 배경 + 주황 테두리)
    add_rect(patent_slide, l, t, w, h, "FFFFFF", line=True)
    # 이미지 (여백 적용)
    patent_slide.shapes.add_picture(
        img_path,
        emu(l + PAD), emu(t + PAD),
        emu(w - PAD*2), emu(h - PAD*2)
    )
    print(f"  ✅ 이미지 {i+1} ({os.path.basename(img_path)}): {l:.2f}\", {t:.2f}\" / {w:.2f}\"x{h:.2f}\"")

# ── 저장 ────────────────────────────────────────────────────────────
prs.save(PPTX_OUT)
print(f"\n✅ 저장 완료: {PPTX_OUT}")

# ── 이 슬라이드만 PNG 재생성 ─────────────────────────────────────────
import subprocess, glob, shutil
LIBRE = r"C:\Program Files\LibreOffice\program\soffice.exe"
TMPDIR = BASE + "/MagicShow강의자료/tmp_page6"
os.makedirs(TMPDIR, exist_ok=True)

print("\n▶ 해당 슬라이드만 PDF export 후 PNG 추출...")

# 전체 PDF 만들기 (페이지 범위 지정이 LibreOffice에서 안 되므로 전체 후 6번만 추출)
result = subprocess.run([
    LIBRE, "--headless", "--convert-to", "pdf",
    "--outdir", TMPDIR, PPTX_OUT
], capture_output=True, text=True, timeout=120)

pdf_files = glob.glob(TMPDIR + "/*.pdf")
if not pdf_files:
    print("ERR: PDF 생성 실패", result.stderr[:200])
else:
    import fitz
    doc = fitz.open(pdf_files[0])
    page = doc[5]  # 0-indexed, slide 6 = index 5
    mat = fitz.Matrix(2.0, 2.0)
    pix = page.get_pixmap(matrix=mat)
    out_png = BASE + "/MagicShow강의자료/slides/slide6.png"
    pix.save(out_png)
    doc.close()
    print(f"✅ slide6.png 교체 완료")

# 임시 파일 정리
shutil.rmtree(TMPDIR)
print("✅ 임시 파일 정리 완료")
